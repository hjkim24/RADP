"""Build ppt/Progress Report_0609.pptx from the 0527 base deck.

Strategy:
  - cp 0527 → 0609 to inherit master/theme (맑은 고딕 + Century Gothic)
  - in-place replace header/page/title text on each slide (positions are
    consistent across the 0527 deck — TEXT_BOX shape[0], shape[2], AUTO_SHAPE
    shape[3])
  - drop any leftover body AUTO_SHAPE / PICTURE shapes (their layouts differ
    per-slide), then ADD a fresh text box for body + add_picture() for the
    figure at a known position
  - trim 0609 to N slides matching SLIDES
"""
from __future__ import annotations
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN

BASE = Path("ppt/Progress Report_0527.pptx")
OUT  = Path("ppt/Progress Report_0609.pptx")
FIGS = Path("paper/figures")

FONT_KO = "맑은 고딕"
FONT_EN = "Times New Roman"

# Slide size = 13.33 × 7.5 inches (16:9). Body region under the title bar.
BODY_LEFT   = Inches(0.5)
BODY_TOP    = Inches(1.7)
BODY_WIDTH  = Inches(12.3)
BODY_HEIGHT = Inches(5.5)

# When a figure is present, body shrinks to the left half
HALF_WIDTH  = Inches(6.6)

SLIDES = [
    {"_kind": "cover"},  # slide 1

    {
        "header": "1. Overview",
        "page": "2",
        "title": "0527 이후 — Progress Overview",
        "body": [
            "▌ 알고리즘",
            "  · α-knob multi-mode (latency / throughput / blended) 도입",
            "  · Async chain forwarding + ResultReady 역방향 채널",
            "  · Mirror cache + Trailer 기반 attribution",
            "",
            "▌ 측정",
            "  · 8-worker Jetson fleet (Orin Nano ×6, AGX Orin MAXN, AGX Xavier)",
            "  · OPT-350M 24-cell matrix + Llama-3.2-1B 9-cell sub-matrix",
            "  · Profiler 두 버그 fix → AGX MAXN 의 16× 우위 처음 노출",
            "",
            "▌ 결과물",
            "  · 영문 LaTeX paper (acmart sigconf, 8p) + 한글 LaTeX paper (KICS, 4p)",
            "  · 모든 section outline → real prose 완료",
        ],
    },

    {
        "header": "2. Algorithm Evolution",
        "page": "3",
        "title": "알고리즘 진화 — 0527 의 단일 throughput DP 에서 확장",
        "body": [
            "▌ α-knob multi-mode",
            "  rank(ψ) = (1−α) Σ T(s) + α · max T(s)",
            "  α=0 latency (EdgeShard) · α=1 throughput (Jupiter) · 0<α<1 blended",
            "  → 같은 DP forward 가 3-mode 를 모두 지원 (cell = (sum, max) tuple)",
            "",
            "▌ Async chain forwarding",
            "  · 동기 chain 의 hidden cost: C streams × N stages = CN 핸들러 thread 점유",
            "  · 비동기 fire-and-forget + ResultReady reverse channel",
            "  · C=16 에서 +17–47% throughput, chain-length-independent",
            "",
            "▌ Mirror cache + Trailer attribution",
            "  · Worker → coordinator 단방향 MirrorActivation RPC (per-stage, ~1.2ms)",
            "  · gRPC trailer 에 실패 stage 의 (start, end) stamp → 원본 ψ_0 에서 lookup",
            "  · Heartbeat 가 substitute 한 후에도 진정으로 죽은 worker 정확 식별",
        ],
    },

    {
        "header": "3. Measurement Infra",
        "page": "4",
        "title": "측정 인프라 + AGX MAXN 의 16× 발견",
        "body": [
            "▌ 8-worker Jetson fleet",
            "  · Orin Nano 8GB × 6",
            "  · AGX Orin (MAXN power mode) × 1",
            "  · AGX Xavier × 1",
            "  · Mac coordinator, 동일 LAN",
            "",
            "▌ 두 profiler 버그 fix",
            "  · 버그 1: tokenizer padding silent no-op",
            "    (seq_length 무관하게 동일 측정)",
            "  · 버그 2: perf_counter() 가 CUDA",
            "    async kernel launch overhead 만 측정",
            "",
            "▌ fix 후 AGX MAXN @ seq=64 decode:",
            "  Nano CUDA 대비 16× faster",
            "  → heterogeneity 가 본격 발현",
        ],
        "figure": FIGS / "fig_tiers.png",
        "figure_pos": (Inches(6.8), Inches(2.0), Inches(6.2), Inches(4.0)),
    },

    {
        "header": "4. Headline Result",
        "page": "5",
        "title": "Headline — L≻T 우위는 두 모델 크기에 걸쳐 C≥4 에서 구조적",
        "body": [
            "▌ 24-cell OPT-350M matrix",
            "  (2 chain × 4 cell × 3 C)",
            "  + 9-cell Llama-1B sub-matrix",
            "",
            "▌ 핵심 finding",
            "  · OPT-350M 24 점 모두에서",
            "    L ≻ T (pairwise)",
            "  · Best: L+async @ C=16 4-stage",
            "    = 41.65 tok/s",
            "  · Worst: T+sync @ C=16 3-stage",
            "    = 23.49 tok/s",
            "  · +77% span — 두 runtime knob",
            "    (placement, chain_mode) 으로 cover",
            "",
            "▌ Llama-1B (3× 큰 모델) 도 동일 pattern",
        ],
        "figure": FIGS / "fig_matrix.png",
        "figure_pos": (Inches(6.4), Inches(1.7), Inches(6.6), Inches(5.5)),
    },

    {
        "header": "5. Async + Model Scaling",
        "page": "6",
        "title": "Async chain + 모델 크기 일반화 — C=1 cross-over 발견",
        "body": [
            "▌ Async > sync, chain-length-independent",
            "  · T-placement @ C=16: 3-stage +47% (23.49 → 34.49 tok/s)",
            "  ·                    4-stage +46% (24.47 → 35.69 tok/s)",
            "  · 동시성에 비례, chain length 에 무관",
            "",
            "▌ Model-size generalisation (Llama-3.2-1B, 3× 큰 모델)",
            "  · L≽T 가 C ≥ 4 에서 유지 (L+async @ C=4 best, +12% vs T+async)",
            "  · 단 C=1 single-stream 에서 T+async > L+async on throughput",
            "    (3.74 vs 3.05 tok/s, −18%)",
            "  · 그러나 TBT p50 는 여전히 L 우위 (174 vs 261 ms, −33%)",
            "",
            "▌ Claim 한정:",
            "  Single-stream 큰 모델 regime 에서 throughput-vs-TBT trade-off 발생",
        ],
    },

    {
        "header": "6. Failure Recovery",
        "page": "7",
        "title": "장애 복구 — Star (5 trials) + Chain (sync/async)",
        "body": [
            "▌ Star topology (RADP)",
            "  · SIGKILL × 5 trials",
            "  · 5/5 trial 60/60 token 전달",
            "  · Recovery: mean 729ms, p95 883ms",
            "",
            "▌ Chain topology (4-stage)",
            "  · mid-chain (on-6) kill",
            "  · Sync: trailer attribution → ~3.6s",
            "  · Async: heartbeat fallback → ~4s",
            "  · ψ+R guarantee 가 correctness 유지",
            "",
            "▌ Baseline 3종 (greedy/uniform/Jupiter-DP)",
            "  · SIGKILL 경계에서 17–20 token 손실",
            "  · stream abort, recovery 없음",
        ],
        "figure": FIGS / "fig_recovery.png",
        "figure_pos": (Inches(6.5), Inches(2.0), Inches(6.5), Inches(4.0)),
    },

    {
        "header": "7. Paper + Backlog",
        "page": "8",
        "title": "논문 작성 + 남은 작업",
        "body": [
            "▌ 논문 작성 (LaTeX 단일 워크플로우)",
            "  · 영문: acmart sigconf, 8p (body 7 + ref 1) — 모든 section prose 완료",
            "  · 한글: KICS 양식, A4 twocolumn, 4p, KoPub Batang/Dotum 폰트",
            "  · 두 양식 동일 references.bib 공유 (Petals/EdgeShard/Jupiter/Helix)",
            "",
            "▌ 단기 (다음 보고 전)",
            "  · Llama-1B L+sync cell 측정 → 4-cell matrix 완성",
            "  · 영문 paper §10.2 baseline 측정 prose 확정",
            "  · 한글 paper 국내 reference 1편 추가 (KICS 권장)",
            "",
            "▌ 중장기",
            "  · A2 — R 단일 백업 → 백업 list 로 확장 (cascading failure)",
            "  · 7B INT4 측정 (Nano memory pressure 해소 후)",
            "  · Cost function 의 marginal-layer term (predicted-vs-measured 2–4% 격차)",
        ],
    },
]


def set_text(tf, text: str, font: str = FONT_KO, size: int = 11, bold: bool = False) -> None:
    """Replace text frame contents and apply font to every new run."""
    lines = text.split("\n")
    tf.text = lines[0]
    for line in lines[1:]:
        p = tf.add_paragraph()
        p.text = line
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.name = font
            run.font.size = Pt(size)
            run.font.bold = bold


def remove_shapes_except(slide, keep_indices: list[int]) -> None:
    """Drop all shapes whose index is not in keep_indices (in-place)."""
    drop = []
    for i, sh in enumerate(slide.shapes):
        if i not in keep_indices:
            drop.append(sh)
    for sh in drop:
        sh._element.getparent().remove(sh._element)


def add_body(slide, lines: list[str], with_figure: bool) -> None:
    """Add a text box for body content under the title bar."""
    left, top = BODY_LEFT, BODY_TOP
    width = HALF_WIDTH if with_figure else BODY_WIDTH
    height = BODY_HEIGHT
    tb = slide.shapes.add_textbox(left, top, width, height)
    set_text(tb.text_frame, "\n".join(lines), FONT_KO, size=12)
    tb.text_frame.word_wrap = True


def main() -> None:
    shutil.copy(BASE, OUT)
    prs = Presentation(OUT)

    # First, trim the deck to len(SLIDES)
    while len(list(prs.slides)) > len(SLIDES):
        xml_slides = prs.slides._sldIdLst
        slides_xml = list(xml_slides)
        xml_slides.remove(slides_xml[-1])

    slides = list(prs.slides)

    for idx, (slide, spec) in enumerate(zip(slides, SLIDES)):
        kind = spec.get("_kind")
        if kind == "cover":
            # Update date + title
            for sh in slide.shapes:
                if not sh.has_text_frame:
                    continue
                t = sh.text_frame.text.strip()
                if t == "Weekly meeting":
                    set_text(sh.text_frame, "2026. 06. 09", FONT_EN, size=24)
                elif t == "RESEARCH PROGRESS":
                    set_text(sh.text_frame, "PROGRESS REPORT", FONT_EN, size=44, bold=True)
            continue

        # Content slide: in-place set header/page/title (shape[0], [2], [3])
        # then strip everything else and rebuild body + figure
        header_shape = slide.shapes[0]
        page_shape   = slide.shapes[2]
        title_shape  = slide.shapes[3]

        set_text(header_shape.text_frame, spec["header"], FONT_EN, size=14, bold=True)
        set_text(page_shape.text_frame, spec["page"], FONT_EN, size=12)
        set_text(title_shape.text_frame, spec["title"], FONT_KO, size=20, bold=True)

        # Drop body/figure shapes (indices 1, 4+) and rebuild
        keep = [0, 2, 3]                      # header, page, title
        # Also keep shape[1] which is just a horizontal-line decoration
        keep.append(1)
        remove_shapes_except(slide, keep)

        with_fig = "figure" in spec
        add_body(slide, spec["body"], with_figure=with_fig)

        if with_fig:
            fig = spec["figure"]
            x, y, w, h = spec["figure_pos"]
            slide.shapes.add_picture(str(fig), x, y, width=w, height=h)

    prs.save(OUT)
    print(f"wrote {OUT}  (slides={len(list(prs.slides))})")


if __name__ == "__main__":
    main()
